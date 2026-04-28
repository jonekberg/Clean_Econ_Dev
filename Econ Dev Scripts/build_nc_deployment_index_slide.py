#!/usr/bin/env python3
"""Build a one-slide PPT ranking North Carolina deployment technologies."""

from __future__ import annotations

import argparse
import os
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_INPUT_CSV = Path(
    os.environ.get(
        "NC_DEPLOYMENT_INPUT_CSV",
        "/Users/jon.ekberg/code/NC-Bootcamp/data/derived/nc_ea_deployment_full.csv",
    )
)
DEFAULT_OUT_DIR = REPO_ROOT / "out" / "nc_deployment_index_slide"

TECH_COLUMNS = {
    "battery_manufacturing": "Battery manufacturing",
    "Solar Generation": "Solar generation",
    "EV_manufacturing": "EV manufacturing",
    "Semiconductor Manufacturing": "Semiconductor manufacturing",
    "solar_manufacturing": "Solar manufacturing",
    "Datacenter": "Data centers",
    "Electricity Storage": "Electricity storage",
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input-csv", type=Path, default=DEFAULT_INPUT_CSV)
    parser.add_argument("--out-dir", type=Path, default=DEFAULT_OUT_DIR)
    return parser.parse_args()


def output_paths(out_dir: Path) -> tuple[Path, Path, Path]:
    return (
        out_dir / "north_carolina_deployment_technology_index.csv",
        out_dir / "north_carolina_deployment_technology_index.png",
        out_dir / "north_carolina_deployment_technology_index_slide.pptx",
    )


def build_index(input_csv: Path) -> pd.DataFrame:
    if not input_csv.exists():
        raise FileNotFoundError(f"Missing input CSV: {input_csv}")

    df = pd.read_csv(input_csv)
    missing = [col for col in TECH_COLUMNS if col not in df.columns]
    if missing:
        raise ValueError(f"Input CSV is missing expected technology columns: {missing}")

    work = df[list(TECH_COLUMNS)].fillna(0)
    totals = work.sum().rename("aggregate_score").reset_index()
    totals.columns = ["technology_raw", "aggregate_score"]
    totals["technology"] = totals["technology_raw"].map(TECH_COLUMNS)
    totals = totals.sort_values("aggregate_score", ascending=False).reset_index(drop=True)
    totals["index_score"] = totals["aggregate_score"] / totals["aggregate_score"].max() * 100
    totals["rank"] = range(1, len(totals) + 1)
    totals["economic_area_count"] = len(df)
    return totals[["rank", "technology", "index_score", "aggregate_score", "economic_area_count"]]


def save_preview(index_df: pd.DataFrame, out_png: Path) -> None:
    plot_df = index_df.sort_values("index_score", ascending=True)
    colors = ["#B9C7DB" for _ in plot_df["technology"]]

    fig, ax = plt.subplots(figsize=(13.333, 7.5), dpi=180)
    fig.patch.set_facecolor("#FFFFFF")
    ax.set_facecolor("#FFFFFF")

    bars = ax.barh(plot_df["technology"], plot_df["index_score"], color=colors, height=0.58)
    ax.set_xlim(0, 100)
    ax.set_xticks([0, 25, 50, 75, 100])
    ax.tick_params(axis="x", labelsize=12, colors="#000000", bottom=False, top=True, labelbottom=False, labeltop=True)
    ax.tick_params(axis="y", labelsize=12, colors="#262626", length=0)
    ax.grid(axis="x", color="#D7DCE3", linewidth=0.8)
    ax.set_axisbelow(True)
    for spine in ax.spines.values():
        spine.set_visible(False)

    top_bar = bars[-1]
    ax.text(
        3,
        top_bar.get_y() + top_bar.get_height() / 2,
        "Battery-led deployment signal",
        va="center",
        ha="left",
        fontsize=12,
        color="#083B63",
        fontweight="bold",
        style="italic",
    )

    fig.text(
        0.075,
        0.91,
        "Leading NC Deployment Categories",
        fontsize=24,
        fontweight="bold",
        color="#000000",
    )
    fig.add_artist(
        plt.Line2D([0.075, 0.16], [0.875, 0.875], transform=fig.transFigure, color="#A9B8D1", linewidth=4)
    )
    fig.text(
        0.075,
        0.82,
        "Index (0-100)",
        fontsize=16,
        color="#000000",
    )
    fig.text(
        0.075,
        0.055,
        "Source: Electrotech_State.R deployment logic; NC economic-area deployment output. Index aggregates technology scores across 23 NC economic areas.",
        fontsize=8.5,
        color="#6B6B6B",
    )
    plt.subplots_adjust(left=0.31, right=0.91, top=0.73, bottom=0.14)
    fig.savefig(out_png, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def add_textbox(slide, left, top, width, height, text, size, bold=False, color="171717"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def disable_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except AttributeError:
        pass


def save_pptx(index_df: pd.DataFrame, out_pptx: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)

    add_textbox(
        slide,
        Inches(0.45),
        Inches(0.48),
        Inches(11.9),
        Inches(0.55),
        "Leading NC Deployment Categories",
        24,
        True,
        "000000",
    )
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.95), Inches(1.25), Inches(0.055))
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor(169, 184, 209)
    underline.line.fill.background()
    disable_shadow(underline)

    add_textbox(slide, Inches(0.46), Inches(1.25), Inches(2.0), Inches(0.34), "Index (0-100)", 17, False, "000000")

    chart_df = index_df.sort_values("index_score", ascending=True)
    chart_left = Inches(3.22)
    chart_top = Inches(1.94)
    chart_width = Inches(8.45)
    chart_height = Inches(4.6)

    chart_data = CategoryChartData()
    chart_data.categories = list(chart_df["technology"])
    chart_data.add_series("Index score", [round(v, 1) for v in chart_df["index_score"]])

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        chart_left,
        chart_top,
        chart_width,
        chart_height,
        chart_data,
    )
    chart = graphic_frame.chart
    chart.has_title = False
    chart.has_legend = False
    chart.plots[0].gap_width = 50

    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(185, 199, 219)
    series.format.line.fill.background()

    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 100
    value_axis.major_unit = 25
    value_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH
    value_axis.tick_labels.font.size = Pt(13)
    value_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(216, 220, 226)
    value_axis.major_gridlines.format.line.width = Pt(0.6)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(12.3)
    category_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)
    category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

    # Keep this callout editable while the bars and workbook remain native chart objects.
    callout = slide.shapes.add_textbox(Inches(6.42), Inches(2.08), Inches(3.25), Inches(0.34))
    tf_bar = callout.text_frame
    tf_bar.clear()
    tf_bar.margin_left = tf_bar.margin_right = tf_bar.margin_top = tf_bar.margin_bottom = 0
    tf_bar.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_bar.paragraphs[0]
    r = p.add_run()
    r.text = "Battery-led deployment signal"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.italic = True
    r.font.color.rgb = RGBColor(8, 59, 99)

    add_textbox(
        slide,
        Inches(0.45),
        Inches(6.88),
        Inches(12.0),
        Inches(0.25),
        "Source: Electrotech_State.R deployment logic; NC economic-area deployment output. Index aggregates technology scores across 23 NC economic areas.",
        8,
        False,
        "77736A",
    )

    prs.save(out_pptx)


def main() -> None:
    args = parse_args()
    args.out_dir.mkdir(parents=True, exist_ok=True)
    out_csv, out_png, out_pptx = output_paths(args.out_dir)
    index_df = build_index(args.input_csv)
    index_df.to_csv(out_csv, index=False)
    save_preview(index_df, out_png)
    save_pptx(index_df, out_pptx)
    print(f"Wrote {out_csv}")
    print(f"Wrote {out_png}")
    print(f"Wrote {out_pptx}")
    print(index_df.to_string(index=False))


if __name__ == "__main__":
    main()
